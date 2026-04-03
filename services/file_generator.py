"""
file_generator.py — Génération de fichiers PDF, Word et PowerPoint
"""
import io
import re


# ── Utilitaires texte ─────────────────────────────────────────────────────────

def sanitize_for_pdf(text):
    """Remplace les caractères Unicode non supportés par Helvetica"""
    replacements = {
        '\u2014': '-', '\u2013': '-', '\u2012': '-', '\u2015': '-',
        '\u2018': "'", '\u2019': "'", '\u201a': ',', '\u201b': "'",
        '\u201c': '"', '\u201d': '"', '\u201e': '"', '\u201f': '"',
        '\u00ab': '"', '\u00bb': '"',
        '\u2026': '...', '\u00a0': ' ', '\u2022': '-', '\u25cf': '-',
        '\u2192': '->', '\u2190': '<-', '\u00b7': '-',
        '\u00d7': 'x', '\u00f7': '/', '\u00ae': '(R)', '\u00a9': '(C)',
        '\u2122': '(TM)', '\u20ac': 'EUR',
        '\u00e0': 'a', '\u00e2': 'a', '\u00e4': 'a', '\u00e7': 'c',
        '\u00e8': 'e', '\u00e9': 'e', '\u00ea': 'e', '\u00eb': 'e',
        '\u00ee': 'i', '\u00ef': 'i', '\u00f4': 'o', '\u00f6': 'o',
        '\u00f9': 'u', '\u00fb': 'u', '\u00fc': 'u', '\u00ff': 'y',
        '\u00c0': 'A', '\u00c2': 'A', '\u00c7': 'C', '\u00c8': 'E',
        '\u00c9': 'E', '\u00ca': 'E', '\u00ce': 'I', '\u00d4': 'O',
        '\u00d9': 'U', '\u00db': 'U', '\u00dc': 'U',
        '\u0153': 'oe', '\u0152': 'OE', '\u00e6': 'ae', '\u00c6': 'AE',
        '\u00df': 'ss', '\u00f1': 'n', '\u00e1': 'a', '\u00e3': 'a',
        '\u00e5': 'a', '\u00ed': 'i', '\u00f3': 'o', '\u00fa': 'u',
        '\u00fd': 'y',
    }
    for char, rep in replacements.items():
        text = text.replace(char, rep)
    return text.encode('ascii', errors='replace').decode('ascii')


def markdown_to_plain(text):
    """Markdown → texte brut ASCII pour PDF"""
    text = re.sub(r'```[\w]*\n?', '', text)
    text = re.sub(r'^#{1,6}\s+(.+)$', r'\1', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    text = re.sub(r'^\s*[-*+\u2022]\s+', '- ', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
    return sanitize_for_pdf(text.strip())


def parse_sections(content):
    """Découpe le contenu markdown en sections (titre, corps)"""
    sections = []
    current_title = None
    current_body = []
    for line in content.split('\n'):
        m = re.match(r'^(#{1,6})\s+(.+)$', line)
        if m:
            if current_body or current_title:
                sections.append((current_title, '\n'.join(current_body).strip()))
            current_title = m.group(2).strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body or current_title:
        sections.append((current_title, '\n'.join(current_body).strip()))
    return sections if sections else [(None, content)]


def strip_markdown(text):
    """Supprime uniquement le markdown de formatage (gras, italique, titres)"""
    text = re.sub(r'#{1,6}\s+', '', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    return text.strip()


# ── PDF ───────────────────────────────────────────────────────────────────────

def generate_pdf(content, title="Document Cognito Chat"):
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            # Ligne fine en haut
            self.set_draw_color(99, 102, 241)
            self.set_line_width(0.3)
            self.line(10, 8, 200, 8)
            # Texte à droite sur la même ligne
            self.set_font('Helvetica', 'I', 8)
            self.set_text_color(150, 150, 170)
            self.set_y(4)
            self.cell(0, 6, 'Cognito Chat - Cognito Inc.', align='R')
            self.ln(8)

        def footer(self):
            self.set_y(-12)
            self.set_draw_color(99, 102, 241)
            self.set_line_width(0.2)
            self.line(10, self.get_y(), 200, self.get_y())
            self.set_font('Helvetica', 'I', 8)
            self.set_text_color(150, 150, 170)
            self.cell(0, 8, f'Page {self.page_no()}', align='C')

    pdf = PDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()
    pdf.set_margins(15, 18, 15)

    # Titre principal
    clean_title = sanitize_for_pdf(strip_markdown(title))
    pdf.set_font('Helvetica', 'B', 20)
    pdf.set_text_color(60, 60, 180)
    pdf.multi_cell(0, 11, clean_title, align='C')
    pdf.ln(4)

    # Ligne de séparation sous le titre
    pdf.set_draw_color(99, 102, 241)
    pdf.set_line_width(0.6)
    pdf.line(15, pdf.get_y(), 195, pdf.get_y())
    pdf.ln(8)

    sections = parse_sections(content)

    for sec_title, sec_body in sections:
        if sec_title:
            pdf.set_font('Helvetica', 'B', 13)
            pdf.set_text_color(40, 40, 120)
            pdf.multi_cell(0, 7, sanitize_for_pdf(strip_markdown(sec_title)))
            pdf.set_draw_color(200, 200, 240)
            pdf.set_line_width(0.2)
            pdf.line(15, pdf.get_y(), 195, pdf.get_y())
            pdf.ln(3)

        pdf.set_font('Helvetica', '', 11)
        pdf.set_text_color(30, 30, 50)

        for line in sec_body.split('\n'):
            line = line.rstrip()
            if not line:
                pdf.ln(2)
                continue
            # Bloc de code
            if line.startswith('    ') or line.startswith('\t'):
                pdf.set_font('Courier', '', 9)
                pdf.set_fill_color(240, 242, 255)
                pdf.multi_cell(0, 5, sanitize_for_pdf(line.strip()), fill=True)
                pdf.set_font('Helvetica', '', 11)
                continue
            # Liste
            if re.match(r'^\s*[-*+\u2022]\s+', line):
                clean = re.sub(r'^\s*[-*+\u2022]\s+', '', line)
                pdf.set_x(20)
                pdf.cell(4, 6, '-')
                pdf.multi_cell(0, 6, markdown_to_plain(clean))
                continue
            # Liste numérotée
            m = re.match(r'^\s*(\d+)\.\s+(.+)$', line)
            if m:
                pdf.set_x(20)
                pdf.cell(6, 6, m.group(1) + '.')
                pdf.multi_cell(0, 6, markdown_to_plain(m.group(2)))
                continue
            # Texte normal
            pdf.multi_cell(0, 6, markdown_to_plain(line))

        pdf.ln(4)

    return bytes(pdf.output())


# ── WORD ──────────────────────────────────────────────────────────────────────

def generate_word(content, title="Document Cognito Chat"):
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # En-tête
    hdr = doc.sections[0].header
    hp = hdr.paragraphs[0]
    hp.text = 'Cognito Chat - Cognito Inc.'
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.runs[0]
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(150, 150, 170)
    run.font.italic = True

    # Titre
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = title_para.add_run(strip_markdown(title))
    tr.font.size = Pt(22)
    tr.font.bold = True
    tr.font.color.rgb = RGBColor(60, 60, 180)
    doc.add_paragraph()

    for sec_title, sec_body in parse_sections(content):
        if sec_title:
            h = doc.add_heading(strip_markdown(sec_title), level=2)
            for run in h.runs:
                run.font.color.rgb = RGBColor(40, 40, 120)

        for line in sec_body.split('\n'):
            line = line.rstrip()
            if not line:
                continue
            if line.startswith('    ') or line.startswith('\t'):
                p = doc.add_paragraph(line.strip(), style='No Spacing')
                p.runs[0].font.name = 'Courier New'
                p.runs[0].font.size = Pt(9)
                continue
            if re.match(r'^\s*[-*+\u2022]\s+', line):
                clean = re.sub(r'^\s*[-*+\u2022]\s+', '', line)
                doc.add_paragraph(strip_markdown(clean), style='List Bullet')
                continue
            if re.match(r'^\s*\d+\.\s+', line):
                clean = re.sub(r'^\s*\d+\.\s+', '', line)
                doc.add_paragraph(strip_markdown(clean), style='List Number')
                continue
            p = doc.add_paragraph()
            parts = re.split(r'\*\*(.+?)\*\*', line)
            if len(parts) > 1:
                for i, part in enumerate(parts):
                    r = p.add_run(part)
                    r.bold = (i % 2 == 1)
            else:
                p.add_run(strip_markdown(line))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── POWERPOINT ────────────────────────────────────────────────────────────────

def generate_pptx(content, title="Présentation Cognito Chat", cover_image=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Palette
    PURPLE  = RGBColor(99,  102, 241)
    CYAN    = RGBColor(6,   182, 212)
    WHITE   = RGBColor(255, 255, 255)
    DARK    = RGBColor(10,  15,  30)
    DARK2   = RGBColor(18,  24,  50)
    DARK3   = RGBColor(25,  35,  70)
    GRAY    = RGBColor(148, 163, 184)
    LIGHT   = RGBColor(226, 232, 240)

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_tb(slide, text, left, top, width, height,
               font_size=18, bold=False, color=WHITE,
               align=PP_ALIGN.LEFT, italic=False, wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = strip_markdown(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txBox

    title_clean = strip_markdown(title)

    # ─── SLIDE TITRE ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)

    has_image = cover_image is not None
    img_start = Inches(7.5)  # image starts at x=7.5" (right 43%)

    if has_image:
        # Image côté droit — ajoutée en premier pour être derrière les formes
        try:
            img_stream = io.BytesIO(cover_image)
            pic = slide.shapes.add_picture(
                img_stream, img_start, 0,
                prs.slide_width - img_start, prs.slide_height
            )
            # Envoyer l'image derrière toutes les autres formes
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)
        except Exception:
            has_image = False

    # Fond gauche opaque (couvre toute la zone texte)
    add_rect(slide, 0, 0, Inches(7.8), prs.slide_height, DARK)

    # Bande décorative droite si pas d'image
    if not has_image:
        add_rect(slide, Inches(8.5), 0, Inches(4.83), prs.slide_height, DARK2)
        add_rect(slide, Inches(10.5), Inches(1), Inches(2.8), Inches(5.5), DARK3)

    # Barre verticale gauche (accent couleur)
    add_rect(slide, 0, 0, Inches(0.12), prs.slide_height, PURPLE)

    # Barre horizontale top
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), PURPLE)

    # Label "COGNITO INC." en haut
    add_tb(slide, 'COGNITO INC.',
           Inches(0.25), Inches(0.18), Inches(5), Inches(0.45),
           font_size=10, bold=True, color=PURPLE)

    # Titre principal
    fs_title = 38 if len(title_clean) <= 35 else (30 if len(title_clean) <= 55 else 24)
    add_tb(slide, title_clean,
           Inches(0.25), Inches(1.9), Inches(7.0), Inches(2.8),
           font_size=fs_title, bold=True, color=WHITE)

    # Ligne décorative sous le titre
    add_rect(slide, Inches(0.25), Inches(4.85), Inches(3.5), Inches(0.05), PURPLE)

    # Sous-titre
    add_tb(slide, 'Cognito Chat - Cognito Inc.',
           Inches(0.25), Inches(5.05), Inches(6.5), Inches(0.55),
           font_size=13, color=GRAY, italic=True)

    # Barre bottom
    add_rect(slide, 0, Inches(7.35), prs.slide_width, Inches(0.15), DARK2)

    # ─── SLIDES DE CONTENU ───────────────────────────────────────────────────
    sections = parse_sections(content)
    slide_num = 0

    for sec_title, sec_body in sections:
        if not sec_title and not sec_body.strip():
            continue
        slide_num += 1

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, DARK)

        # Bande header
        add_rect(slide, 0, 0, prs.slide_width, Inches(1.25), DARK2)
        # Accent gauche violet
        add_rect(slide, 0, 0, Inches(0.1), Inches(1.25), PURPLE)
        # Accent droit cyan
        add_rect(slide, prs.slide_width - Inches(0.1), 0, Inches(0.1), prs.slide_height, CYAN)

        # Numéro de slide
        add_tb(slide, f'{slide_num:02d}',
               Inches(12.5), Inches(0.1), Inches(0.7), Inches(0.55),
               font_size=13, bold=True, color=PURPLE, align=PP_ALIGN.RIGHT)

        # Titre section
        display_title = strip_markdown(sec_title) if sec_title else 'Contenu'
        add_tb(slide, display_title,
               Inches(0.25), Inches(0.08), Inches(11.8), Inches(1.1),
               font_size=27, bold=True, color=WHITE)

        # Ligne de séparation sous header
        add_rect(slide, Inches(0.1), Inches(1.25), Inches(13.1), Inches(0.03), PURPLE)

        # Extraire les bullets
        bullets = []
        for line in sec_body.split('\n'):
            line = line.rstrip()
            if not line:
                continue
            clean = re.sub(r'^\s*[-*+\u2022\d]+[.)]\s*', '', line).strip()
            clean = strip_markdown(clean)
            if clean and len(clean) > 2:
                bullets.append(clean)

        if bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.4), Inches(1.45), Inches(12.7), Inches(5.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            max_bullets = 6
            fs_bullet = 20 if len(bullets) <= 4 else (17 if len(bullets) <= 6 else 15)

            for i, bullet in enumerate(bullets[:max_bullets]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Emu(120000)
                p.space_after  = Emu(40000)
                run = p.add_run()
                run.text = '\u25b8  ' + bullet   # ▸
                run.font.size = Pt(fs_bullet)
                run.font.color.rgb = LIGHT if i % 2 == 0 else GRAY

        # Footer
        add_rect(slide, 0, Inches(7.2), prs.slide_width, Inches(0.3), DARK2)
        add_tb(slide, 'Cognito Chat - Cognito Inc.',
               Inches(0), Inches(7.2), Inches(13.33), Inches(0.3),
               font_size=8, color=RGBColor(80, 95, 130),
               align=PP_ALIGN.CENTER, italic=True)

    # ─── SLIDE FINALE ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)

    # Bloc décoratif droit
    add_rect(slide, Inches(8.5), 0,        Inches(4.83), prs.slide_height, DARK2)
    add_rect(slide, Inches(10),  Inches(1.2), Inches(3.33), Inches(5.1),   DARK3)
    add_rect(slide, Inches(11.5), Inches(2.5), Inches(1.83), Inches(2.5),  PURPLE)

    # Accents
    add_rect(slide, 0, 0, Inches(0.12), prs.slide_height, PURPLE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), PURPLE)
    add_rect(slide, 0, Inches(7.44), prs.slide_width, Inches(0.06), PURPLE)

    add_tb(slide, 'Merci',
           Inches(0.25), Inches(2.2), Inches(8), Inches(1.4),
           font_size=60, bold=True, color=WHITE)

    add_rect(slide, Inches(0.25), Inches(3.75), Inches(4), Inches(0.06), PURPLE)

    add_tb(slide, title_clean,
           Inches(0.25), Inches(3.95), Inches(7.8), Inches(0.9),
           font_size=18, color=CYAN)

    add_tb(slide, 'Cognito Chat - Cognito Inc.',
           Inches(0.25), Inches(5.0), Inches(7.8), Inches(0.55),
           font_size=13, color=GRAY, italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Détection de demande de fichier ──────────────────────────────────────────

def should_generate_file(text):
    """
    Retourne ('pdf'|'word'|'pptx', sujet) ou (None, None)
    """
    text_lower = text.lower()

    action_words = [
        'génère', 'genere', 'générer', 'generer', 'crée', 'cree', 'créer', 'creer',
        'fais', 'faire', 'produis', 'produire', 'rédige', 'redige', 'écris', 'ecris',
        'donne', 'donner', 'fais-moi', 'fais moi', 'donne-moi', 'donne moi',
        'generate', 'create', 'make', 'write', 'produce', 'give me'
    ]
    has_action = any(w in text_lower for w in action_words)

    pptx_words = ['powerpoint', 'pptx', '.ppt', 'présentation ppt', 'presentation ppt',
                  'diaporama', 'slideshow', 'slides', 'présentation powerpoint',
                  'presentation powerpoint']
    pdf_words  = ['pdf', 'fichier pdf', 'document pdf', 'en pdf', 'au format pdf', '.pdf']
    word_words = ['word', 'docx', 'fichier word', 'document word', 'en word',
                  'au format word', '.docx', 'microsoft word']

    is_pptx = any(w in text_lower for w in pptx_words)
    is_pdf  = any(w in text_lower for w in pdf_words)
    is_word = any(w in text_lower for w in word_words)

    if not (is_pptx or is_pdf or is_word):
        return None, None

    # Extraire le sujet en retirant les mots de commande et de format
    remove_words = [
        'génère', 'genere', 'crée', 'cree', 'fais moi', 'fais-moi', 'donne moi', 'donne-moi',
        'un pdf', 'une pdf', 'un word', 'une word', 'un fichier', 'un document',
        'une presentation', 'une présentation', 'un powerpoint', 'un diaporama',
        'au format pdf', 'en pdf', 'au format word', 'en word', 'au format ppt',
        'en ppt', 'en powerpoint', '.pdf', '.docx', '.pptx', '.ppt',
        'pdf', 'word', 'docx', 'powerpoint', 'pptx', 'diaporama',
        'avec des images', 'avec images',
    ]
    subject = text
    for w in sorted(remove_words, key=len, reverse=True):
        subject = re.sub(re.escape(w), ' ', subject, flags=re.IGNORECASE)
    subject = re.sub(r'\s+', ' ', subject).strip(' ,.-:;?!')
    subject = strip_markdown(subject)
    if not subject or len(subject) < 3:
        subject = text

    if is_pptx:
        return 'pptx', subject
    if is_pdf:
        return 'pdf', subject
    return 'word', subject
